"""
Report generation service for PDF and PPTX exports.
"""

import io
import logging
from datetime import datetime, timedelta, timezone

from sqlalchemy import func
from sqlalchemy.orm import Session

from app.models.station import MonitoringStation
from app.models.timeseries import TimeSeriesData
from app.models.alert import Alert
from app.services.aqi_calculator import calculate_aqi, get_aqi_category

logger = logging.getLogger(__name__)


def generate_report(
    db: Session,
    station_ids: list[str],
    start_date: str,
    end_date: str,
    report_format: str = "pdf",
    sections: list[str] = None,
) -> io.BytesIO:
    """
    Generate a PDF or PPTX environmental report.
    """
    if sections is None:
        sections = ["aqi", "water", "noise", "weather"]

    # Gather data
    report_data = _gather_report_data(db, station_ids, start_date, end_date)

    if report_format == "pptx":
        return _generate_pptx(report_data, sections)
    else:
        return _generate_pdf(report_data, sections)


def _gather_report_data(
    db: Session,
    station_ids: list[str],
    start_date: str,
    end_date: str,
) -> dict:
    """Gather all data needed for the report."""
    start = datetime.fromisoformat(start_date)
    end = datetime.fromisoformat(end_date)

    stations_data = []
    for sid in station_ids:
        station = db.query(MonitoringStation).filter(MonitoringStation.station_id == sid).first()
        if not station:
            continue

        # Get aggregated data
        averages = (
            db.query(
                func.avg(TimeSeriesData.aqi_pm25).label("avg_pm25"),
                func.avg(TimeSeriesData.aqi_pm10).label("avg_pm10"),
                func.avg(TimeSeriesData.aqi_so2).label("avg_so2"),
                func.avg(TimeSeriesData.aqi_no2).label("avg_no2"),
                func.avg(TimeSeriesData.aqi_co).label("avg_co"),
                func.avg(TimeSeriesData.aqi_o3).label("avg_o3"),
                func.avg(TimeSeriesData.water_ph).label("avg_ph"),
                func.avg(TimeSeriesData.water_bod).label("avg_bod"),
                func.avg(TimeSeriesData.water_cod).label("avg_cod"),
                func.avg(TimeSeriesData.water_tss).label("avg_tss"),
                func.avg(TimeSeriesData.noise_leq).label("avg_noise"),
                func.avg(TimeSeriesData.noise_lday).label("avg_lday"),
                func.avg(TimeSeriesData.noise_lnight).label("avg_lnight"),
                func.avg(TimeSeriesData.weather_temp).label("avg_temp"),
                func.avg(TimeSeriesData.weather_humidity).label("avg_humidity"),
                func.max(TimeSeriesData.aqi_pm25).label("max_pm25"),
                func.min(TimeSeriesData.aqi_pm25).label("min_pm25"),
            )
            .filter(
                TimeSeriesData.station_id == sid,
                TimeSeriesData.timestamp >= start,
                TimeSeriesData.timestamp <= end,
            )
            .first()
        )

        # Count alerts
        alert_count = (
            db.query(func.count(Alert.alert_id))
            .filter(
                Alert.station_id == sid,
                Alert.timestamp >= start,
                Alert.timestamp <= end,
            )
            .scalar()
        ) or 0

        aqi_result = calculate_aqi(
            pm25=averages.avg_pm25 if averages else None,
            pm10=averages.avg_pm10 if averages else None,
            so2=averages.avg_so2 if averages else None,
            no2=averages.avg_no2 if averages else None,
            co=averages.avg_co if averages else None,
            o3=averages.avg_o3 if averages else None,
        )

        stations_data.append({
            "station": {
                "id": station.station_id,
                "name": station.station_name,
                "city": station.city,
                "state": station.state,
                "zone_type": station.zone_type,
            },
            "averages": {
                "pm25": round(averages.avg_pm25, 1) if averages and averages.avg_pm25 else None,
                "pm10": round(averages.avg_pm10, 1) if averages and averages.avg_pm10 else None,
                "so2": round(averages.avg_so2, 1) if averages and averages.avg_so2 else None,
                "no2": round(averages.avg_no2, 1) if averages and averages.avg_no2 else None,
                "co": round(averages.avg_co, 2) if averages and averages.avg_co else None,
                "o3": round(averages.avg_o3, 1) if averages and averages.avg_o3 else None,
                "noise": round(averages.avg_noise, 1) if averages and averages.avg_noise else None,
                "ph": round(averages.avg_ph, 2) if averages and averages.avg_ph else None,
                "bod": round(averages.avg_bod, 1) if averages and averages.avg_bod else None,
                "temp": round(averages.avg_temp, 1) if averages and averages.avg_temp else None,
            },
            "aqi": aqi_result,
            "alerts": alert_count,
            "max_pm25": round(averages.max_pm25, 1) if averages and averages.max_pm25 else None,
            "min_pm25": round(averages.min_pm25, 1) if averages and averages.min_pm25 else None,
        })

    return {
        "stations": stations_data,
        "period": {"start": start_date, "end": end_date},
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def _generate_pdf(report_data: dict, sections: list[str]) -> io.BytesIO:
    """Generate PDF report using ReportLab."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    )

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    elements = []

    # Title style
    title_style = ParagraphStyle(
        "ReportTitle", parent=styles["Heading1"], fontSize=24, spaceAfter=20,
        textColor=colors.HexColor("#1a5276"), alignment=1
    )
    subtitle_style = ParagraphStyle(
        "Subtitle", parent=styles["Normal"], fontSize=12, spaceAfter=10,
        textColor=colors.HexColor("#566573"), alignment=1
    )
    heading2_style = ParagraphStyle(
        "Heading2Custom", parent=styles["Heading2"], fontSize=16, spaceAfter=10,
        textColor=colors.HexColor("#1a5276")
    )

    # Cover page
    elements.append(Spacer(1, 2 * inch))
    elements.append(Paragraph("PrithviNet", title_style))
    elements.append(Paragraph("Environmental Monitoring Report", subtitle_style))
    elements.append(Spacer(1, 0.5 * inch))
    period = report_data["period"]
    elements.append(Paragraph(f"Period: {period['start']} to {period['end']}", subtitle_style))
    elements.append(Paragraph(f"Generated: {report_data['generated_at'][:10]}", subtitle_style))
    elements.append(Paragraph(
        f"Stations: {len(report_data['stations'])}", subtitle_style
    ))
    elements.append(PageBreak())

    # Per-station reports
    for sdata in report_data["stations"]:
        station = sdata["station"]
        avgs = sdata["averages"]
        aqi = sdata["aqi"]

        elements.append(Paragraph(f"{station['name']}", heading2_style))
        elements.append(Paragraph(
            f"{station['city']}, {station['state']} | Zone: {station['zone_type']} | Alerts: {sdata['alerts']}",
            styles["Normal"]
        ))
        elements.append(Spacer(1, 0.3 * inch))

        # AQI Section
        if "aqi" in sections and aqi.get("aqi"):
            elements.append(Paragraph("Air Quality Index", styles["Heading3"]))
            aqi_data = [
                ["Metric", "Value"],
                ["Overall AQI", str(aqi.get("aqi", "N/A"))],
                ["Category", aqi.get("category", "N/A")],
                ["Prominent Pollutant", aqi.get("prominent_pollutant", "N/A")],
                ["Avg PM2.5 (ug/m3)", str(avgs.get("pm25", "N/A"))],
                ["Avg PM10 (ug/m3)", str(avgs.get("pm10", "N/A"))],
                ["Avg SO2 (ug/m3)", str(avgs.get("so2", "N/A"))],
                ["Avg NO2 (ug/m3)", str(avgs.get("no2", "N/A"))],
                ["Avg CO (mg/m3)", str(avgs.get("co", "N/A"))],
                ["Avg O3 (ug/m3)", str(avgs.get("o3", "N/A"))],
            ]
            t = Table(aqi_data, colWidths=[3 * inch, 3 * inch])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a5276")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f2f4f4")]),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 0.3 * inch))

        # Water Section
        if "water" in sections and avgs.get("ph"):
            elements.append(Paragraph("Water Quality", styles["Heading3"]))
            water_data = [
                ["Metric", "Value"],
                ["pH", str(avgs.get("ph", "N/A"))],
                ["BOD (mg/L)", str(avgs.get("bod", "N/A"))],
            ]
            t = Table(water_data, colWidths=[3 * inch, 3 * inch])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a5276")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 0.3 * inch))

        # Noise Section
        if "noise" in sections and avgs.get("noise"):
            elements.append(Paragraph("Noise Levels", styles["Heading3"]))
            noise_data = [
                ["Metric", "Value"],
                ["Leq (dB)", str(avgs.get("noise", "N/A"))],
            ]
            t = Table(noise_data, colWidths=[3 * inch, 3 * inch])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a5276")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 0.3 * inch))

        elements.append(PageBreak())

    doc.build(elements)
    buffer.seek(0)
    return buffer


def _generate_pptx(report_data: dict, sections: list[str]) -> io.BytesIO:
    """Generate PPTX report using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "PrithviNet"
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Environmental Monitoring Report"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x56, 0x65, 0x73)
    p.alignment = PP_ALIGN.CENTER

    period = report_data["period"]
    p = tf.add_paragraph()
    p.text = f"Period: {period['start']} to {period['end']}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0x85, 0x92, 0x9E)
    p.alignment = PP_ALIGN.CENTER

    # Per-station slides
    for sdata in report_data["stations"]:
        station = sdata["station"]
        avgs = sdata["averages"]
        aqi = sdata["aqi"]

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Station title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = station["name"]
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = f"{station['city']}, {station['state']} | {station['zone_type']} Zone | {sdata['alerts']} Alerts"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x56, 0x65, 0x73)

        # Metrics grid
        y_pos = Inches(1.8)
        metrics = []
        if "aqi" in sections and aqi.get("aqi"):
            metrics.append(("AQI", str(aqi["aqi"]), aqi.get("category", "")))
        if "aqi" in sections and avgs.get("pm25"):
            metrics.append(("PM2.5", f"{avgs['pm25']} ug/m3", ""))
        if "noise" in sections and avgs.get("noise"):
            metrics.append(("Noise", f"{avgs['noise']} dB", ""))
        if "water" in sections and avgs.get("ph"):
            metrics.append(("Water pH", str(avgs["ph"]), ""))
        if "weather" in sections and avgs.get("temp"):
            metrics.append(("Temperature", f"{avgs['temp']} C", ""))

        for i, (label, value, subtext) in enumerate(metrics):
            x = Inches(0.5 + (i % 4) * 3.2)
            y = y_pos + Inches((i // 4) * 2.2)

            txBox = slide.shapes.add_textbox(x, y, Inches(2.8), Inches(1.8))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x56, 0x65, 0x73)

            p = tf.add_paragraph()
            p.text = value
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)

            if subtext:
                p = tf.add_paragraph()
                p.text = subtext
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0x85, 0x92, 0x9E)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
